from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0D, 0x0D, 0x0D)   # near-black
ACCENT     = RGBColor(0xD4, 0x7A, 0x21)   # Claude orange
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GREY   = RGBColor(0x88, 0x88, 0x88)
TAG_BG     = RGBColor(0x1E, 0x1E, 0x1E)   # slightly lighter than BG
TAG_TEXT   = ACCENT

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helpers ─────────────────────────────────────────────────────────

def bg(slide, color=BG_DARK):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def txbox(slide, text, left, top, width, height,
          size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          wrap=True):
    """Add a simple text box."""
    box = slide.shapes.add_textbox(left, top, width, height)
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box

def accent_bar(slide, top=Inches(0.18), height=Inches(0.06)):
    """Thin orange bar at top."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, top, SLIDE_W, height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

def tag(slide, label, left, top, width=Inches(2.4), height=Inches(0.38)):
    """Small pill tag (e.g. SCREENSHARE / SLIDES)."""
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = TAG_BG
    box.line.color.rgb = ACCENT
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = TAG_TEXT

def bullets(slide, items, left, top, width, height,
            size=20, color=LIGHT_GREY, spacing_after=Pt(10)):
    """Multi-bullet text box with em-bullet character."""
    box = slide.shapes.add_textbox(left, top, width, height)
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing_after
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color

def divider(slide, top, color=ACCENT, alpha=80):
    """Horizontal rule."""
    line = slide.shapes.add_shape(1, Inches(0.6), top, Inches(12.1), Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txbox(s, "NEXTLEAP  ·  APPLIED GENAI BOOTCAMP",
      Inches(0.6), Inches(0.7), Inches(12), Inches(0.5),
      size=13, color=MID_GREY)

txbox(s, "Getting Hands On\nwith Claude Code",
      Inches(0.6), Inches(1.4), Inches(10), Inches(2.4),
      size=52, bold=True, color=WHITE)

txbox(s, "Session 1 of 2   ·   March 21, 2025   ·   Aman Parmar",
      Inches(0.6), Inches(4.0), Inches(10), Inches(0.5),
      size=16, color=MID_GREY)

divider(s, Inches(3.7))

# bottom tagline
txbox(s, "By the end of this session: you'll have Claude Code running,\nyour project context configured, and fear of the terminal — gone.",
      Inches(0.6), Inches(4.6), Inches(11), Inches(1.2),
      size=19, color=LIGHT_GREY)


# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda / Time Breakdown
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txbox(s, "What We're Covering Today", Inches(0.6), Inches(0.5),
      Inches(10), Inches(0.7), size=32, bold=True)

rows = [
    ("0:00 – 0:15", "Why Claude Code for PMs"),
    ("0:15 – 0:25", "What is Claude Code"),
    ("0:25 – 0:45", "Live Setup Walkthrough"),
    ("0:45 – 1:10", "CLAUDE.md — The Brain"),
    ("1:10 – 1:30", "Memory File Hierarchy"),
    ("1:30 – 1:45", "Project File Structure"),
    ("1:45 – 1:55", "Live Showcase — End to End"),
    ("1:55 – 2:00", "Q&A + Build Day Setup"),
]

row_h = Inches(0.58)
top_start = Inches(1.35)
for i, (time, title) in enumerate(rows):
    top = top_start + i * row_h
    # alternating row shade
    if i % 2 == 0:
        shade = s.shapes.add_shape(1, Inches(0.5), top + Inches(0.04),
                                   Inches(12.3), row_h - Inches(0.06))
        shade.fill.solid()
        shade.fill.fore_color.rgb = TAG_BG
        shade.line.fill.background()

    txbox(s, time,  Inches(0.65), top + Inches(0.1), Inches(1.7),  row_h, size=14, color=ACCENT, bold=True)
    txbox(s, title, Inches(2.45), top + Inches(0.1), Inches(10.0), row_h, size=16, color=WHITE)


# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — You already know AI / Hook
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txbox(s, "You Already Know AI.\nThis Session Is About Using It Like a Builder.",
      Inches(0.6), Inches(0.55), Inches(11), Inches(1.9),
      size=34, bold=True)
divider(s, Inches(2.6))

pts = [
    "Week 8 means you know prompts, agents, RAG, and evals.",
    "Most PMs still copy-paste into a chat window.",
    "Claude Code is the step between 'user' and 'builder'.",
    "The terminal is not scary — it's just a faster keyboard.",
]
bullets(s, pts, Inches(0.6), Inches(2.75), Inches(10.5), Inches(3.5), size=22)

txbox(s, "Today's goal: run Claude Code yourself before this session ends.",
      Inches(0.6), Inches(6.3), Inches(11), Inches(0.7),
      size=18, bold=True, color=ACCENT)


# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — Claude Code is a coding agent
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txbox(s, "Claude Code Is a Coding Agent, Not a Chat Tool",
      Inches(0.6), Inches(0.55), Inches(11.5), Inches(1.1),
      size=34, bold=True)
divider(s, Inches(1.75))

pts = [
    "Lives in your terminal — not a browser tab.",
    "Reads your files, writes code, runs commands.",
    "You don't describe what you want — you give it context and a task.",
    "ChatGPT answers questions.  Claude Code executes work.",
]
bullets(s, pts, Inches(0.6), Inches(2.0), Inches(7.5), Inches(3.5), size=22)

# right panel — comparison box
box = s.shapes.add_shape(1, Inches(8.6), Inches(1.9), Inches(4.2), Inches(4.8))
box.fill.solid()
box.fill.fore_color.rgb = TAG_BG
box.line.color.rgb = ACCENT

txbox(s, "ChatGPT / Gemini", Inches(8.75), Inches(2.1), Inches(3.8), Inches(0.5),
      size=14, bold=True, color=ACCENT)
txbox(s, "→ Answers in a window\n→ Forgets after session\n→ Can't touch your files",
      Inches(8.75), Inches(2.55), Inches(3.8), Inches(1.4), size=14, color=LIGHT_GREY)

txbox(s, "Claude Code", Inches(8.75), Inches(4.05), Inches(3.8), Inches(0.5),
      size=14, bold=True, color=ACCENT)
txbox(s, "→ Runs in your project\n→ Reads real files\n→ Executes and iterates",
      Inches(8.75), Inches(4.5), Inches(3.8), Inches(1.4), size=14, color=LIGHT_GREY)


# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — When to use Claude Code
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txbox(s, "When PMs Should Reach for Claude Code",
      Inches(0.6), Inches(0.55), Inches(11.5), Inches(1.0),
      size=34, bold=True)
divider(s, Inches(1.65))

use_cases = [
    ("Building with a real file structure",
     "Not a one-off answer — an actual project with folders, configs, scripts."),
    ("Working inside an existing repo",
     "Claude reads the codebase. Context is automatic. No copy-pasting."),
    ("Needing Claude to run and test things",
     "It executes commands, checks output, iterates. Not just suggests."),
    ("Repeatable workflows",
     "PRD reviews, standup notes, bug triage — teach it once, reuse forever."),
]

card_w = Inches(2.9)
card_h = Inches(4.2)
card_top = Inches(1.9)
for i, (title, desc) in enumerate(use_cases):
    left = Inches(0.4) + i * (card_w + Inches(0.28))
    card = s.shapes.add_shape(1, left, card_top, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = TAG_BG
    card.line.color.rgb = ACCENT

    txbox(s, title, left + Inches(0.15), card_top + Inches(0.2),
          card_w - Inches(0.3), Inches(0.85), size=16, bold=True, color=ACCENT)
    txbox(s, desc,  left + Inches(0.15), card_top + Inches(1.1),
          card_w - Inches(0.3), Inches(2.8), size=15, color=LIGHT_GREY)


# ════════════════════════════════════════════════════════════════════
# SLIDE 6 — Three things that make it powerful
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txbox(s, "Three Things That Make Claude Code Powerful",
      Inches(0.6), Inches(0.55), Inches(11.5), Inches(1.0),
      size=34, bold=True)
divider(s, Inches(1.65))

txbox(s, "Today covers #1 and #2. Session 2 covers #3.",
      Inches(0.6), Inches(1.75), Inches(11), Inches(0.5),
      size=16, color=MID_GREY)

layers = [
    ("01", "CLAUDE.md",     "Persistent context. Tells Claude who you are, how you work, what the project is.", Inches(1.8)),
    ("02", "Memory Files",  "Structured memory that persists across sessions. Claude picks up where you left off.", Inches(3.4)),
    ("03", "Skills",        "Repeatable workflows you teach Claude once and reuse forever. Session 2.", Inches(5.0)),
]

for num, name, desc, top in layers:
    bar = s.shapes.add_shape(1, Inches(0.6), top, Inches(12.1), Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TAG_BG
    bar.line.color.rgb = ACCENT

    txbox(s, num,  Inches(0.8),  top + Inches(0.2), Inches(0.6),  Inches(1.0), size=24, bold=True, color=ACCENT)
    txbox(s, name, Inches(1.55), top + Inches(0.15), Inches(2.5), Inches(0.55), size=20, bold=True, color=WHITE)
    txbox(s, desc, Inches(1.55), top + Inches(0.7),  Inches(9.5), Inches(0.55), size=15, color=LIGHT_GREY)


# ════════════════════════════════════════════════════════════════════
# SLIDE 7 — SCREENSHARE: Setup Walkthrough
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txbox(s, "Live Setup Walkthrough",
      Inches(0.6), Inches(0.55), Inches(10), Inches(1.0),
      size=36, bold=True)
divider(s, Inches(1.65))

steps = [
    "npm install -g @anthropic-ai/claude-code  (or via brew)",
    "claude — first launch, API key setup",
    "claude --help — show it's just a CLI",
    "Run a basic task: 'List all files and tell me what this project does'",
    "Show /clear, /help, /exit",
]
bullets(s, steps, Inches(0.6), Inches(1.85), Inches(10.5), Inches(3.8), size=21)

txbox(s, "The terminal fear is real. The cure is one command that works.",
      Inches(0.6), Inches(6.3), Inches(11), Inches(0.7),
      size=17, bold=True, color=ACCENT)


# ════════════════════════════════════════════════════════════════════
# SLIDE 8 — CLAUDE.md is the brain
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txbox(s, "CLAUDE.md Is the Context File\nClaude Reads Before Every Session",
      Inches(0.6), Inches(0.5), Inches(9.5), Inches(1.8),
      size=34, bold=True)
divider(s, Inches(2.45))

pts = [
    "Sits at the root of your project.",
    "Claude reads it automatically on startup — every time.",
    "Think of it as your project brief that never disappears.",
    "Without it, Claude has no memory of what you're building.",
]
bullets(s, pts, Inches(0.6), Inches(2.65), Inches(7.5), Inches(3.2), size=21)

# right panel
panel = s.shapes.add_shape(1, Inches(8.5), Inches(2.45), Inches(4.3), Inches(4.7))
panel.fill.solid()
panel.fill.fore_color.rgb = TAG_BG
panel.line.color.rgb = ACCENT

txbox(s, "project-root/\n  ├── CLAUDE.md  ← Claude reads this first\n  ├── src/\n  ├── scripts/\n  └── docs/",
      Inches(8.65), Inches(2.7), Inches(3.9), Inches(2.0),
      size=13, color=ACCENT)

txbox(s, "Every session starts here.\nNo setup. No re-explaining.\nJust context.",
      Inches(8.65), Inches(4.9), Inches(3.9), Inches(1.5),
      size=14, color=LIGHT_GREY)


# ════════════════════════════════════════════════════════════════════
# SLIDE 9 — What goes in CLAUDE.md
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txbox(s, "What Goes in CLAUDE.md",
      Inches(0.6), Inches(0.55), Inches(10), Inches(0.9),
      size=34, bold=True)
divider(s, Inches(1.55))

sections = [
    ("Project Identity",  "What is this? What problem does it solve? Who is it for?"),
    ("Tech Stack",        "Language, frameworks, conventions, version constraints."),
    ("File Structure",    "Where things live and why. What not to touch."),
    ("Rules",             "What Claude should never do. Hard constraints."),
    ("Quick Reference",   "Commands, URLs, env setup, anything Claude needs fast."),
]

row_h = Inches(0.88)
for i, (title, desc) in enumerate(sections):
    top = Inches(1.75) + i * row_h
    num_box = s.shapes.add_shape(1, Inches(0.5), top, Inches(0.45), Inches(0.6))
    num_box.fill.solid()
    num_box.fill.fore_color.rgb = ACCENT
    num_box.line.fill.background()
    txbox(s, str(i+1), Inches(0.52), top, Inches(0.4), Inches(0.6),
          size=15, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

    txbox(s, title, Inches(1.1), top + Inches(0.02), Inches(3.2), Inches(0.55),
          size=17, bold=True, color=WHITE)
    txbox(s, desc,  Inches(4.5), top + Inches(0.02), Inches(8.2), Inches(0.55),
          size=16, color=LIGHT_GREY)

    if i < len(sections) - 1:
        divider(s, top + Inches(0.76), color=RGBColor(0x33, 0x33, 0x33))


# ════════════════════════════════════════════════════════════════════
# SLIDE 10 — What NOT to put in CLAUDE.md
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txbox(s, "What NOT to Put in CLAUDE.md",
      Inches(0.6), Inches(0.55), Inches(10), Inches(0.9),
      size=34, bold=True)
divider(s, Inches(1.55))

bad  = ["Credentials or API keys", "Content that changes every session", "Long essays — Claude reads the whole file every time", "Instructions for every edge case"]
good = ["Store secrets in .env", "Use memory files for session-state", "Keep it scannable. One line per rule.", "Cover the 80%. Trust Claude for the rest."]

txbox(s, "INSTEAD", Inches(7.3), Inches(1.65), Inches(3), Inches(0.45),
      size=13, bold=True, color=ACCENT)
txbox(s, "AVOID",   Inches(0.6), Inches(1.65), Inches(3), Inches(0.45),
      size=13, bold=True, color=RGBColor(0xFF, 0x55, 0x55))

for i, (b, g) in enumerate(zip(bad, good)):
    top = Inches(2.2) + i * Inches(1.0)
    bad_box = s.shapes.add_shape(1, Inches(0.5), top, Inches(6.3), Inches(0.75))
    bad_box.fill.solid()
    bad_box.fill.fore_color.rgb = RGBColor(0x2A, 0x10, 0x10)
    bad_box.line.color.rgb = RGBColor(0xFF, 0x55, 0x55)
    txbox(s, f"✗  {b}", Inches(0.65), top + Inches(0.08), Inches(6.0), Inches(0.6),
          size=16, color=RGBColor(0xFF, 0x88, 0x88))

    good_box = s.shapes.add_shape(1, Inches(7.0), top, Inches(5.8), Inches(0.75))
    good_box.fill.solid()
    good_box.fill.fore_color.rgb = RGBColor(0x0D, 0x2A, 0x10)
    good_box.line.color.rgb = RGBColor(0x44, 0xBB, 0x44)
    txbox(s, f"✓  {g}", Inches(7.15), top + Inches(0.08), Inches(5.5), Inches(0.6),
          size=16, color=RGBColor(0x88, 0xDD, 0x88))


# ════════════════════════════════════════════════════════════════════
# SLIDE 11 — SCREENSHARE: CLAUDE.md Live Demo
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txbox(s, "CLAUDE.md — Live Demo",
      Inches(0.6), Inches(0.55), Inches(10), Inches(1.0),
      size=36, bold=True)
divider(s, Inches(1.65))

steps = [
    "Open a real project — walk through the CLAUDE.md as-is",
    "Start a new Claude session — show it picks up context immediately",
    "Run the same task with CLAUDE.md vs without — show the quality gap",
    "Live edit: add one new rule, restart, confirm Claude follows it",
]
bullets(s, steps, Inches(0.6), Inches(1.85), Inches(10.5), Inches(3.4), size=21)

txbox(s, "10 minutes of setup = weeks of better output.",
      Inches(0.6), Inches(6.3), Inches(11), Inches(0.7),
      size=17, bold=True, color=ACCENT)


# ════════════════════════════════════════════════════════════════════
# SLIDE 12 — Memory hierarchy
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txbox(s, "Memory File Hierarchy — Three Levels",
      Inches(0.6), Inches(0.55), Inches(11), Inches(0.9),
      size=34, bold=True)
divider(s, Inches(1.55))

levels = [
    ("GLOBAL",    "~/.claude/CLAUDE.md",      "Every Claude session",        "Your preferences, writing rules, personal tools"),
    ("PROJECT",   "project-root/CLAUDE.md",   "Any session in that folder",  "Project context, stack, conventions, hard rules"),
    ("SUBFOLDER", "subfolder/CLAUDE.md",       "Sessions inside that folder", "Module overrides, component-specific rules"),
]

col_w = Inches(4.0)
for i, (label, path, scope, use) in enumerate(levels):
    left = Inches(0.4) + i * (col_w + Inches(0.16))
    card = s.shapes.add_shape(1, left, Inches(1.75), col_w, Inches(4.9))
    card.fill.solid()
    card.fill.fore_color.rgb = TAG_BG
    card.line.color.rgb = ACCENT

    txbox(s, label, left + Inches(0.15), Inches(1.95), col_w - Inches(0.3), Inches(0.5),
          size=13, bold=True, color=ACCENT)
    txbox(s, path,  left + Inches(0.15), Inches(2.5),  col_w - Inches(0.3), Inches(0.55),
          size=14, bold=True, color=WHITE)
    txbox(s, f"When: {scope}", left + Inches(0.15), Inches(3.15), col_w - Inches(0.3), Inches(0.55),
          size=13, color=MID_GREY)
    txbox(s, use, left + Inches(0.15), Inches(3.8), col_w - Inches(0.3), Inches(2.4),
          size=14, color=LIGHT_GREY)

txbox(s, "Claude reads all three. More specific = higher priority.",
      Inches(0.6), Inches(6.85), Inches(11), Inches(0.5),
      size=16, bold=True, color=ACCENT)


# ════════════════════════════════════════════════════════════════════
# SLIDE 13 — Global vs Project split
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txbox(s, "The Right Split — What Lives Where",
      Inches(0.6), Inches(0.55), Inches(11), Inches(0.9),
      size=34, bold=True)
divider(s, Inches(1.55))

cols = [
    ("Global\n~/.claude/CLAUDE.md",    ["How you like code structured", "Your writing voice and tone rules", "Tools you use across every project", "Personal defaults and shortcuts"]),
    ("Project\nproject/CLAUDE.md",      ["What this product does", "File conventions for this repo", "Hard rules: never do X in this codebase", "Stack and framework versions"]),
    ("Memory Files\n.claude/memory/",   ["In-progress work and open decisions", "Decisions made this sprint", "What to pick up next session", "Running notes and context"]),
]

col_w = Inches(3.9)
for i, (header, items) in enumerate(cols):
    left = Inches(0.4) + i * (col_w + Inches(0.2))
    card = s.shapes.add_shape(1, left, Inches(1.75), col_w, Inches(5.0))
    card.fill.solid()
    card.fill.fore_color.rgb = TAG_BG
    card.line.color.rgb = ACCENT

    txbox(s, header, left + Inches(0.15), Inches(1.9), col_w - Inches(0.3), Inches(0.9),
          size=16, bold=True, color=WHITE)
    bullets(s, items, left + Inches(0.1), Inches(2.9), col_w - Inches(0.2), Inches(3.5),
            size=15, color=LIGHT_GREY, spacing_after=Pt(8))


# ════════════════════════════════════════════════════════════════════
# SLIDE 14 — Project File Structure
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txbox(s, "Recommended Project Structure for PM Projects",
      Inches(0.6), Inches(0.55), Inches(11), Inches(1.0),
      size=33, bold=True)
divider(s, Inches(1.65))

# code block
code_bg = s.shapes.add_shape(1, Inches(0.5), Inches(1.85), Inches(5.8), Inches(5.2))
code_bg.fill.solid()
code_bg.fill.fore_color.rgb = RGBColor(0x0A, 0x0A, 0x0A)
code_bg.line.color.rgb = ACCENT

code = (
    "project-name/\n"
    "├── CLAUDE.md          ← project brain\n"
    "├── .claude/\n"
    "│   ├── memory/        ← session state\n"
    "│   └── skills/        ← reusable workflows\n"
    "├── data/              ← inputs, raw files\n"
    "├── scripts/           ← automation scripts\n"
    "├── outputs/           ← generated content\n"
    "└── docs/              ← specs, decisions"
)
txbox(s, code, Inches(0.7), Inches(2.0), Inches(5.4), Inches(4.8),
      size=14, color=ACCENT)

txbox(s, "Works for:", Inches(6.8), Inches(1.95), Inches(5.8), Inches(0.45),
      size=16, bold=True, color=WHITE)
use_list = [
    "Content systems",
    "Research and synthesis tools",
    "Reporting and analytics scripts",
    "Internal dashboards",
    "Anything a PM builds with Claude",
]
bullets(s, use_list, Inches(6.8), Inches(2.45), Inches(5.8), Inches(4.2),
        size=17, color=LIGHT_GREY)


# ════════════════════════════════════════════════════════════════════
# SLIDE 15 — Live Showcase
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txbox(s, "Live Showcase — End to End",
      Inches(0.6), Inches(0.55), Inches(10), Inches(1.0),
      size=36, bold=True)
divider(s, Inches(1.65))

steps = [
    ("1", "New project folder — zero setup"),
    ("2", "Create CLAUDE.md live (2 minutes, from template)"),
    ("3", "Real PM task: 'Read this PRD and generate a test plan'"),
    ("4", "Claude reads the file, writes the output, saves it"),
    ("5", "Review the output — explain what made it good"),
]

for num, step in steps:
    top = Inches(1.9) + (int(num)-1) * Inches(0.8)
    nb = s.shapes.add_shape(1, Inches(0.5), top, Inches(0.5), Inches(0.55))
    nb.fill.solid()
    nb.fill.fore_color.rgb = ACCENT
    nb.line.fill.background()
    txbox(s, num, Inches(0.52), top, Inches(0.45), Inches(0.55),
          size=15, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    txbox(s, step, Inches(1.15), top + Inches(0.04), Inches(10.5), Inches(0.55),
          size=18, color=WHITE)

txbox(s, "12 minutes from zero to working output. Setup cost is front-loaded. Every session after is faster.",
      Inches(0.6), Inches(6.7), Inches(11), Inches(0.6),
      size=16, bold=True, color=ACCENT)


# ════════════════════════════════════════════════════════════════════
# SLIDE 16 — Sunday Build Day starters
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txbox(s, "Three Starter Tasks for Sunday's Build Day",
      Inches(0.6), Inches(0.55), Inches(11), Inches(0.9),
      size=34, bold=True)

txbox(s, "Pick one. Finish it. Ship something.",
      Inches(0.6), Inches(1.45), Inches(8), Inches(0.5),
      size=18, color=MID_GREY)
divider(s, Inches(2.05))

tasks = [
    ("01", "Create Your CLAUDE.md",
     "For an existing project or a new idea.\nUse the template from today as your starting point."),
    ("02", "Automate a Repeatable Task",
     "PRD review, standup notes, bug triage.\nWrite the CLAUDE.md, give Claude the task, iterate until it's good."),
    ("03", "Scaffold a New Project",
     "Give Claude the goal, let it propose the folder structure.\nRefine it. You own the output."),
]

card_w = Inches(3.9)
for i, (num, title, desc) in enumerate(tasks):
    left = Inches(0.4) + i * (card_w + Inches(0.2))
    card = s.shapes.add_shape(1, left, Inches(2.25), card_w, Inches(4.7))
    card.fill.solid()
    card.fill.fore_color.rgb = TAG_BG
    card.line.color.rgb = ACCENT

    txbox(s, num, left + Inches(0.15), Inches(2.4), card_w - Inches(0.3), Inches(0.55),
          size=28, bold=True, color=ACCENT)
    txbox(s, title, left + Inches(0.15), Inches(3.0), card_w - Inches(0.3), Inches(0.7),
          size=18, bold=True, color=WHITE)
    txbox(s, desc, left + Inches(0.15), Inches(3.75), card_w - Inches(0.3), Inches(2.8),
          size=15, color=LIGHT_GREY)


# ════════════════════════════════════════════════════════════════════
# SLIDE 17 — What's in Session 2 (teaser)
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txbox(s, "Today you gave Claude a brain.\nSession 2: you teach it habits.",
      Inches(0.6), Inches(0.6), Inches(11), Inches(1.9),
      size=38, bold=True)
divider(s, Inches(2.6))

txbox(s, "Session 2 — Agent Skills", Inches(0.6), Inches(2.8), Inches(6), Inches(0.7),
      size=24, bold=True, color=ACCENT)

pts = [
    "What are skills — and why they're different from CLAUDE.md",
    "Skill levels: metadata, instructions, resources",
    "Skills architecture — how Claude loads and uses them",
    "Build a skill from scratch — live, in session",
    "PM skill showcase — real workflows automated",
]
bullets(s, pts, Inches(0.6), Inches(3.55), Inches(11), Inches(3.0), size=20)


# ════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════
out = "/Users/amanparmar/Documents/AI-PM/Projects/ship-with-ai/next-leap/Session-1-Claude-Code-NextLeap.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
