"""Layout eval for ppt-builder decks.

Run after saving a .pptx to catch the two layout bugs that have shipped most often:
  1. Shapes that extend past the slide canvas (off-slide content).
  2. Shapes that overlap when they should not (footer over last row, label over title, etc.).

Usage:
    python3 check_layout.py <path-to-pptx>

Exit code:
    0 if no issues found
    1 if any issue found (so it can gate a CI / pre-delivery step)

What it checks:
    - Out-of-bounds: any shape whose bounding box leaves the 13.33 x 7.5 inch canvas
    - Overlap: any pair of text-bearing shapes whose bounding boxes overlap by more than
      a small tolerance. Decorative shapes (the orange accent bar at the top of every
      slide, full-width dividers, background rectangles) are excluded by heuristic.

What it does NOT check:
    - Text overflowing inside a shape (would require font measurement)
    - Color contrast, font size minima, accessibility
    - Slide count or content correctness

Tune thresholds via the constants near the top of the file.
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

SLIDE_W_EMU = Emu(13.33 * 914400)
SLIDE_H_EMU = Emu(7.5 * 914400)

OVERLAP_TOLERANCE_EMU = Emu(0.05 * 914400)
DECORATIVE_HEIGHT_EMU = Emu(0.12 * 914400)
FULL_WIDTH_THRESHOLD = 0.95


def shape_bbox(shape):
    return (shape.left or 0, shape.top or 0,
            (shape.left or 0) + (shape.width or 0),
            (shape.top or 0) + (shape.height or 0))


def is_decorative(shape) -> bool:
    """Skip the orange accent bar, full-width dividers, and slide-fill background rectangles."""
    if shape.width is None or shape.height is None:
        return True
    width_ratio = shape.width / SLIDE_W_EMU
    height_ratio = shape.height / SLIDE_H_EMU
    if width_ratio >= FULL_WIDTH_THRESHOLD and shape.height <= DECORATIVE_HEIGHT_EMU:
        return True
    if width_ratio >= 0.99 and height_ratio >= 0.99:
        return True
    return False


def has_text(shape) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    return any(run.text.strip()
               for para in shape.text_frame.paragraphs
               for run in para.runs)


def rectangles_overlap(a, b, tol):
    ax1, ay1, ax2, ay2 = a
    bx1, by1, bx2, by2 = b
    if ax2 <= bx1 + tol or bx2 <= ax1 + tol:
        return False
    if ay2 <= by1 + tol or by2 <= ay1 + tol:
        return False
    return True


def emu_to_in(value):
    return round(value / 914400, 2)


def fmt_bbox(bbox):
    return (f"({emu_to_in(bbox[0])}, {emu_to_in(bbox[1])}) -> "
            f"({emu_to_in(bbox[2])}, {emu_to_in(bbox[3])}) inches")


def first_text(shape, n=50):
    if not getattr(shape, "has_text_frame", False):
        return ""
    txt = shape.text_frame.text.strip().replace("\n", " ")
    return (txt[:n] + "...") if len(txt) > n else txt


def check_slide(idx, slide):
    issues = []
    shapes = list(slide.shapes)

    for sh in shapes:
        if sh.left is None or sh.top is None:
            continue
        bbox = shape_bbox(sh)
        if (bbox[2] > SLIDE_W_EMU + OVERLAP_TOLERANCE_EMU
                or bbox[3] > SLIDE_H_EMU + OVERLAP_TOLERANCE_EMU
                or bbox[0] < -OVERLAP_TOLERANCE_EMU
                or bbox[1] < -OVERLAP_TOLERANCE_EMU):
            issues.append(
                f"  OUT OF BOUNDS: {fmt_bbox(bbox)} | text: '{first_text(sh)}'"
            )

    text_shapes = [sh for sh in shapes
                   if has_text(sh) and not is_decorative(sh)
                   and sh.left is not None and sh.top is not None]
    for i, a in enumerate(text_shapes):
        for b in text_shapes[i + 1:]:
            if rectangles_overlap(shape_bbox(a), shape_bbox(b), OVERLAP_TOLERANCE_EMU):
                issues.append(
                    f"  OVERLAP: '{first_text(a)}'  <->  '{first_text(b)}'"
                )

    return issues


def main():
    if len(sys.argv) != 2:
        print("Usage: python3 check_layout.py <path-to-pptx>")
        sys.exit(2)
    path = Path(sys.argv[1])
    if not path.exists():
        print(f"File not found: {path}")
        sys.exit(2)

    prs = Presentation(str(path))
    total_issues = 0
    for idx, slide in enumerate(prs.slides, 1):
        issues = check_slide(idx, slide)
        if issues:
            print(f"Slide {idx}:")
            for line in issues:
                print(line)
            total_issues += len(issues)

    if total_issues:
        print(f"\nFAIL: {total_issues} layout issue(s) found in {path.name}")
        sys.exit(1)
    print(f"OK: no layout issues in {path.name} ({len(prs.slides)} slides)")
    sys.exit(0)


if __name__ == "__main__":
    main()
