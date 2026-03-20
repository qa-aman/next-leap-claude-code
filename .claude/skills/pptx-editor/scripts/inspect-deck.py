"""
Inspect a PPTX deck and extract its design tokens.

This is Phase 1 of the pptx-editor skill workflow.
Run this before making any edits to understand the deck's design system.

Usage:
    python3 .claude/skills/pptx-editor/scripts/inspect-deck.py <path-to-pptx>

Example:
    python3 .claude/skills/pptx-editor/scripts/inspect-deck.py 02-presentation/deck.pptx

Output:
    Prints design tokens (colors, fonts, positions, shapes) for every slide.
    Use this output to match the deck's visual style when adding or editing slides.
"""

import sys
from pptx import Presentation
from pptx.oxml.ns import qn


def get_fill_color(fill):
    """Extract fill color safely."""
    if fill.type is None:
        return None
    try:
        return str(fill.fore_color.rgb)
    except Exception:
        return "(complex/theme)"


def get_font_color(font):
    """Extract font color safely."""
    try:
        if font.color and font.color.rgb:
            return str(font.color.rgb)
    except Exception:
        pass
    return "theme"


def get_typeface(run):
    """Extract font typeface - explicit name or THEME."""
    rPr = run._r.find(qn('a:rPr'))
    if rPr is None:
        return "THEME"
    latin = rPr.find(qn('a:latin'))
    if latin is None:
        return "THEME"
    return latin.get('typeface', 'THEME')


def inspect(path):
    prs = Presentation(path)

    print("=" * 70)
    print("DECK OVERVIEW")
    print("=" * 70)
    print(f"File: {path}")
    print(f"Slide width:  {prs.slide_width} EMU")
    print(f"Slide height: {prs.slide_height} EMU")
    print(f"Total slides: {len(prs.slides)}")
    print()

    # Collect design tokens across slides
    bg_colors = set()
    accent_colors = set()
    text_styles = []

    for idx, slide in enumerate(prs.slides):
        bg_fill = slide.background.fill
        bg_color = get_fill_color(bg_fill) if bg_fill.type else "theme/none"
        bg_colors.add(bg_color)

        print(f"{'=' * 70}")
        print(f"SLIDE {idx + 1} | Layout: {slide.slide_layout.name}")
        print(f"{'=' * 70}")
        print(f"  Background: type={bg_fill.type}, color={bg_color}")
        print()

        for shape in slide.shapes:
            shape_fill = get_fill_color(shape.fill) if hasattr(shape, 'fill') else None

            # Shape info
            print(f"  [{shape.name}]")
            print(f"    pos=({shape.left}, {shape.top})  size=({shape.width}, {shape.height})")

            if shape_fill:
                print(f"    fill: {shape_fill}")
                accent_colors.add(shape_fill)

            # Line info
            if shape.line and shape.line.fill and shape.line.fill.type is not None:
                print(f"    line: type={shape.line.fill.type}, width={shape.line.width}")

            # Text info
            if shape.has_text_frame:
                for pi, para in enumerate(shape.text_frame.paragraphs):
                    if not para.text.strip():
                        continue
                    if para.runs:
                        r = para.runs[0]
                        typeface = get_typeface(r)
                        color = get_font_color(r.font)
                        text_preview = para.text[:80].replace('\n', ' ')
                        style_info = {
                            "size": r.font.size,
                            "bold": r.font.bold,
                            "color": color,
                            "font": typeface,
                        }
                        text_styles.append(style_info)
                        print(f"    text[{pi}]: '{text_preview}'")
                        print(f"      size={r.font.size}  bold={r.font.bold}  color={color}  font={typeface}")
                    else:
                        print(f"    text[{pi}]: '{para.text[:80]}' (no runs)")

            print()

    # Summary
    print("=" * 70)
    print("DESIGN TOKEN SUMMARY")
    print("=" * 70)
    print(f"Background colors: {bg_colors}")
    print(f"Shape fill colors: {accent_colors}")
    print()

    # Find common text styles
    sizes = {}
    for s in text_styles:
        key = f"size={s['size']} bold={s['bold']} color={s['color']} font={s['font']}"
        sizes[key] = sizes.get(key, 0) + 1

    print("Text styles (by frequency):")
    for style, count in sorted(sizes.items(), key=lambda x: -x[1]):
        print(f"  {count}x  {style}")


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print(f"Usage: python3 {sys.argv[0]} <path-to-pptx>")
        sys.exit(1)
    inspect(sys.argv[1])
