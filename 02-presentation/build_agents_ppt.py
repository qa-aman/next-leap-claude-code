from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette ──────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0D, 0x0D, 0x0D)
ACCENT     = RGBColor(0xD4, 0x7A, 0x21)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GREY   = RGBColor(0x88, 0x88, 0x88)
TAG_BG     = RGBColor(0x1E, 0x1E, 0x1E)
TAG_TEXT   = ACCENT

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ── Helpers ─────────────────────────────────────────────────────────
def bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def txbox(slide, text, left, top, width, height,
          size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box

def accent_bar(slide, top=Inches(0.18), height=Inches(0.06)):
    bar = slide.shapes.add_shape(1, 0, top, SLIDE_W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

def bullets(slide, items, left, top, width, height,
            size=20, color=LIGHT_GREY, spacing_after=Pt(10)):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing_after
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color

def divider(slide, top, color=ACCENT):
    line = slide.shapes.add_shape(1, Inches(0.6), top, Inches(12.1), Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "BASED ON ANTHROPIC'S ENGINEERING GUIDANCE",
      Inches(0.6), Inches(0.7), Inches(12), Inches(0.5),
      size=13, color=MID_GREY)

txbox(s, "Building Effective\nAI Agents",
      Inches(0.6), Inches(1.4), Inches(11), Inches(2.6),
      size=52, bold=True, color=WHITE)

txbox(s, "What agents are, when to use them, and the patterns that work",
      Inches(0.6), Inches(4.0), Inches(12), Inches(0.5),
      size=16, color=MID_GREY)

divider(s, Inches(4.7))

txbox(s, "Source: anthropic.com/engineering/building-effective-agents",
      Inches(0.6), Inches(5.0), Inches(11), Inches(0.7),
      size=16, color=LIGHT_GREY)


# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "What We'll Cover", Inches(0.6), Inches(0.5),
      Inches(10), Inches(0.7), size=32, bold=True)

rows = [
    ("01", "Agents vs Workflows — the core distinction"),
    ("02", "When to use each — and when to use neither"),
    ("03", "Five workflow patterns Anthropic recommends"),
    ("04", "How agents actually work — the autonomous loop"),
    ("05", "Core principles for building reliable agents"),
    ("06", "Common pitfalls and what to do instead"),
]

row_h = Inches(0.7)
top_start = Inches(1.45)
for i, (num, title) in enumerate(rows):
    top = top_start + i * row_h
    if i % 2 == 0:
        shade = s.shapes.add_shape(1, Inches(0.5), top + Inches(0.04),
                                   Inches(12.3), row_h - Inches(0.06))
        shade.fill.solid()
        shade.fill.fore_color.rgb = TAG_BG
        shade.line.fill.background()
    txbox(s, num,   Inches(0.75), top + Inches(0.12), Inches(1.0), row_h, size=16, color=ACCENT, bold=True)
    txbox(s, title, Inches(2.05), top + Inches(0.12), Inches(10.5), row_h, size=18, color=WHITE)


# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — Start Simple Hook
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "Before You Build an Agent —\nAsk if You Actually Need One.",
      Inches(0.6), Inches(0.55), Inches(12), Inches(2.0),
      size=34, bold=True)
divider(s, Inches(2.7))

pts = [
    "Anthropic's first guidance: start simple.",
    "Optimize a single LLM call with retrieval and in-context examples first.",
    "Add complexity only when simpler solutions fall short.",
    "Agents trade latency and cost for flexibility — make sure that trade is worth it.",
]
bullets(s, pts, Inches(0.6), Inches(2.95), Inches(11), Inches(3.3), size=22)

txbox(s, "\"For many applications, a single optimized LLM call is enough.\"",
      Inches(0.6), Inches(6.3), Inches(12), Inches(0.7),
      size=17, bold=True, color=ACCENT)


# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — Agents vs Workflows
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "Agents vs Workflows — Anthropic's Definition",
      Inches(0.6), Inches(0.55), Inches(12), Inches(1.0),
      size=32, bold=True)
divider(s, Inches(1.65))

# Workflows card
left = Inches(0.5)
card_w = Inches(6.1)
card_h = Inches(5.0)
top = Inches(1.95)

card = s.shapes.add_shape(1, left, top, card_w, card_h)
card.fill.solid(); card.fill.fore_color.rgb = TAG_BG
card.line.color.rgb = ACCENT
txbox(s, "WORKFLOWS", left + Inches(0.2), top + Inches(0.25),
      card_w - Inches(0.4), Inches(0.5), size=14, bold=True, color=ACCENT)
txbox(s, "Predefined code paths",
      left + Inches(0.2), top + Inches(0.75),
      card_w - Inches(0.4), Inches(0.6), size=22, bold=True, color=WHITE)
txbox(s,
      "LLMs and tools are orchestrated through code you write. "
      "You decide the steps. The model fills in the blanks.\n\n"
      "Best for: well-defined tasks where predictability and "
      "consistency matter more than flexibility.",
      left + Inches(0.2), top + Inches(1.55),
      card_w - Inches(0.4), Inches(3.2), size=16, color=LIGHT_GREY)

# Agents card
left2 = Inches(6.85)
card2 = s.shapes.add_shape(1, left2, top, card_w, card_h)
card2.fill.solid(); card2.fill.fore_color.rgb = TAG_BG
card2.line.color.rgb = ACCENT
txbox(s, "AGENTS", left2 + Inches(0.2), top + Inches(0.25),
      card_w - Inches(0.4), Inches(0.5), size=14, bold=True, color=ACCENT)
txbox(s, "Model-directed processes",
      left2 + Inches(0.2), top + Inches(0.75),
      card_w - Inches(0.4), Inches(0.6), size=22, bold=True, color=WHITE)
txbox(s,
      "The LLM dynamically directs its own steps and tool use. "
      "It maintains control over how to accomplish the task.\n\n"
      "Best for: open-ended problems where you can't predict "
      "the number of steps in advance.",
      left2 + Inches(0.2), top + Inches(1.55),
      card_w - Inches(0.4), Inches(3.2), size=16, color=LIGHT_GREY)


# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — When to use which
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "Choose the Pattern That Fits the Task",
      Inches(0.6), Inches(0.55), Inches(12), Inches(1.0),
      size=32, bold=True)
divider(s, Inches(1.65))

cols = [
    ("Use a Single LLM Call", [
        "Task is well-scoped",
        "One round of reasoning is enough",
        "Latency and cost matter most",
        "You can hand-craft the prompt",
    ]),
    ("Use a Workflow", [
        "Steps are knowable in advance",
        "You need consistency across runs",
        "Quality gates between steps help",
        "Predictability beats flexibility",
    ]),
    ("Use an Agent", [
        "Open-ended, exploratory task",
        "Step count can't be predicted",
        "Model needs to react to tool output",
        "Flexibility justifies the cost",
    ]),
]

col_w = Inches(4.0)
for i, (header, items) in enumerate(cols):
    left = Inches(0.4) + i * (col_w + Inches(0.16))
    card = s.shapes.add_shape(1, left, Inches(1.95), col_w, Inches(5.0))
    card.fill.solid(); card.fill.fore_color.rgb = TAG_BG
    card.line.color.rgb = ACCENT
    txbox(s, header, left + Inches(0.2), Inches(2.1),
          col_w - Inches(0.4), Inches(0.7), size=18, bold=True, color=WHITE)
    bullets(s, items, left + Inches(0.15), Inches(2.95),
            col_w - Inches(0.3), Inches(3.8), size=15, color=LIGHT_GREY, spacing_after=Pt(10))


# ════════════════════════════════════════════════════════════════════
# SLIDE 6 — Five Workflow Patterns Overview
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "Five Workflow Patterns Anthropic Recommends",
      Inches(0.6), Inches(0.55), Inches(12), Inches(1.0),
      size=32, bold=True)
divider(s, Inches(1.65))

layers = [
    ("01", "Prompt Chaining",     "Decompose a task into sequential steps with quality gates between them.", Inches(1.85)),
    ("02", "Routing",             "Classify inputs and direct each one to a specialized downstream handler.",  Inches(2.85)),
    ("03", "Parallelization",     "Run multiple LLM calls simultaneously via sectioning or voting.",          Inches(3.85)),
    ("04", "Orchestrator-Workers","A central LLM breaks tasks down and delegates to worker models.",          Inches(4.85)),
    ("05", "Evaluator-Optimizer", "One LLM generates output; another evaluates and gives iterative feedback.", Inches(5.85)),
]

for num, name, desc, top in layers:
    bar = s.shapes.add_shape(1, Inches(0.6), top, Inches(12.1), Inches(0.85))
    bar.fill.solid(); bar.fill.fore_color.rgb = TAG_BG
    bar.line.color.rgb = ACCENT
    txbox(s, num,  Inches(0.8),  top + Inches(0.15), Inches(0.7),  Inches(0.6), size=20, bold=True, color=ACCENT)
    txbox(s, name, Inches(1.55), top + Inches(0.1),  Inches(3.2),  Inches(0.5), size=17, bold=True, color=WHITE)
    txbox(s, desc, Inches(4.85), top + Inches(0.13), Inches(7.7),  Inches(0.5), size=14, color=LIGHT_GREY)


# ════════════════════════════════════════════════════════════════════
# SLIDE 7 — Prompt Chaining
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "Pattern 1 — Prompt Chaining",
      Inches(0.6), Inches(0.55), Inches(12), Inches(1.0),
      size=32, bold=True)
divider(s, Inches(1.65))

pts = [
    "Break the task into a sequence of LLM calls.",
    "Each step's output becomes the next step's input.",
    "Add programmatic checks (gates) between steps to catch errors early.",
    "Trades a little latency for much higher accuracy.",
]
bullets(s, pts, Inches(0.6), Inches(1.9), Inches(7.5), Inches(3.5), size=20)

panel = s.shapes.add_shape(1, Inches(8.5), Inches(1.85), Inches(4.3), Inches(4.4))
panel.fill.solid(); panel.fill.fore_color.rgb = TAG_BG
panel.line.color.rgb = ACCENT

txbox(s, "EXAMPLE FLOW", Inches(8.65), Inches(2.0),
      Inches(4.0), Inches(0.4), size=13, bold=True, color=ACCENT)
txbox(s,
      "Outline →\n   Gate: outline valid?\n\n"
      "Draft →\n   Gate: covers all sections?\n\n"
      "Polish →\n   Final output",
      Inches(8.65), Inches(2.5), Inches(4.0), Inches(3.5),
      size=15, color=LIGHT_GREY)

txbox(s, "Good for: writing, translation, multi-step analysis.",
      Inches(0.6), Inches(6.5), Inches(11), Inches(0.6),
      size=16, bold=True, color=ACCENT)


# ════════════════════════════════════════════════════════════════════
# SLIDE 8 — Routing
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "Pattern 2 — Routing",
      Inches(0.6), Inches(0.55), Inches(12), Inches(1.0),
      size=32, bold=True)
divider(s, Inches(1.65))

pts = [
    "A classifier LLM reads the input and decides where to send it.",
    "Each downstream handler is specialized for one input type.",
    "Lets you use the right model and prompt for each case.",
    "Separates concerns — easier to optimize each branch independently.",
]
bullets(s, pts, Inches(0.6), Inches(1.9), Inches(7.5), Inches(3.5), size=20)

panel = s.shapes.add_shape(1, Inches(8.5), Inches(1.85), Inches(4.3), Inches(4.4))
panel.fill.solid(); panel.fill.fore_color.rgb = TAG_BG
panel.line.color.rgb = ACCENT

txbox(s, "EXAMPLE FLOW", Inches(8.65), Inches(2.0),
      Inches(4.0), Inches(0.4), size=13, bold=True, color=ACCENT)
txbox(s,
      "Customer query →\n   Classify intent\n\n"
      "  • Refund → refund agent\n"
      "  • Bug    → support agent\n"
      "  • Sales  → sales agent\n\n"
      "Right handler, every time.",
      Inches(8.65), Inches(2.5), Inches(4.0), Inches(3.5),
      size=14, color=LIGHT_GREY)

txbox(s, "Good for: support triage, query classification, mixed-intent inputs.",
      Inches(0.6), Inches(6.5), Inches(11), Inches(0.6),
      size=16, bold=True, color=ACCENT)


# ════════════════════════════════════════════════════════════════════
# SLIDE 9 — Parallelization
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "Pattern 3 — Parallelization",
      Inches(0.6), Inches(0.55), Inches(12), Inches(1.0),
      size=32, bold=True)
divider(s, Inches(1.65))

pts = [
    "Run multiple LLM calls at the same time, then combine results.",
    "Two flavors: sectioning (split work) and voting (consensus).",
    "Sectioning cuts latency on independent subtasks.",
    "Voting improves reliability by aggregating multiple attempts.",
]
bullets(s, pts, Inches(0.6), Inches(1.9), Inches(7.5), Inches(3.5), size=20)

panel = s.shapes.add_shape(1, Inches(8.5), Inches(1.85), Inches(4.3), Inches(4.4))
panel.fill.solid(); panel.fill.fore_color.rgb = TAG_BG
panel.line.color.rgb = ACCENT

txbox(s, "TWO FLAVORS", Inches(8.65), Inches(2.0),
      Inches(4.0), Inches(0.4), size=13, bold=True, color=ACCENT)
txbox(s,
      "Sectioning\n"
      "Split a doc into chunks,\nsummarize each in parallel.\n\n"
      "Voting\n"
      "Generate 5 answers,\npick the most common one.",
      Inches(8.65), Inches(2.5), Inches(4.0), Inches(3.5),
      size=14, color=LIGHT_GREY)

txbox(s, "Good for: speed-up on independent work, higher-confidence answers.",
      Inches(0.6), Inches(6.5), Inches(11), Inches(0.6),
      size=16, bold=True, color=ACCENT)


# ════════════════════════════════════════════════════════════════════
# SLIDE 10 — Orchestrator-Workers
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "Pattern 4 — Orchestrator-Workers",
      Inches(0.6), Inches(0.55), Inches(12), Inches(1.0),
      size=32, bold=True)
divider(s, Inches(1.65))

pts = [
    "A central LLM (the orchestrator) plans and delegates work.",
    "Worker LLMs do focused sub-tasks and report back.",
    "The orchestrator decides what to do next based on results.",
    "More flexible than parallelization — the plan adapts at runtime.",
]
bullets(s, pts, Inches(0.6), Inches(1.9), Inches(7.5), Inches(3.5), size=20)

panel = s.shapes.add_shape(1, Inches(8.5), Inches(1.85), Inches(4.3), Inches(4.4))
panel.fill.solid(); panel.fill.fore_color.rgb = TAG_BG
panel.line.color.rgb = ACCENT

txbox(s, "EXAMPLE FLOW", Inches(8.65), Inches(2.0),
      Inches(4.0), Inches(0.4), size=13, bold=True, color=ACCENT)
txbox(s,
      "Orchestrator: 'Refactor\nthis codebase.'\n\n"
      "  → Worker 1: rename vars\n"
      "  → Worker 2: extract funcs\n"
      "  → Worker 3: update tests\n\n"
      "Orchestrator plans next round.",
      Inches(8.65), Inches(2.5), Inches(4.0), Inches(3.5),
      size=13, color=LIGHT_GREY)

txbox(s, "Good for: coding tasks, research, anything where the plan emerges.",
      Inches(0.6), Inches(6.5), Inches(11), Inches(0.6),
      size=16, bold=True, color=ACCENT)


# ════════════════════════════════════════════════════════════════════
# SLIDE 11 — Evaluator-Optimizer
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "Pattern 5 — Evaluator-Optimizer",
      Inches(0.6), Inches(0.55), Inches(12), Inches(1.0),
      size=32, bold=True)
divider(s, Inches(1.65))

pts = [
    "One LLM generates a response. Another evaluates and gives feedback.",
    "The generator revises. The loop continues until the evaluator is satisfied.",
    "Mimics how humans iterate on drafts.",
    "Best when you can articulate clear evaluation criteria.",
]
bullets(s, pts, Inches(0.6), Inches(1.9), Inches(7.5), Inches(3.5), size=20)

panel = s.shapes.add_shape(1, Inches(8.5), Inches(1.85), Inches(4.3), Inches(4.4))
panel.fill.solid(); panel.fill.fore_color.rgb = TAG_BG
panel.line.color.rgb = ACCENT

txbox(s, "EXAMPLE FLOW", Inches(8.65), Inches(2.0),
      Inches(4.0), Inches(0.4), size=13, bold=True, color=ACCENT)
txbox(s,
      "Generator: write a story.\n\n"
      "Evaluator: 'Tone is off.\nMake it darker.'\n\n"
      "Generator: revise.\n\n"
      "Evaluator: 'Approved.'",
      Inches(8.65), Inches(2.6), Inches(4.0), Inches(3.4),
      size=14, color=LIGHT_GREY)

txbox(s, "Good for: writing, code review, translation, anything with clear quality criteria.",
      Inches(0.6), Inches(6.5), Inches(12), Inches(0.6),
      size=16, bold=True, color=ACCENT)


# ════════════════════════════════════════════════════════════════════
# SLIDE 12 — How an Agent Actually Works
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "How an Agent Actually Works — The Loop",
      Inches(0.6), Inches(0.55), Inches(12), Inches(1.0),
      size=32, bold=True)
divider(s, Inches(1.65))

steps = [
    ("1", "Receive the task and any starting context"),
    ("2", "Plan the next step based on the goal"),
    ("3", "Call a tool — run code, search, read a file"),
    ("4", "Observe the result and update its plan"),
    ("5", "Repeat until the task is done or a stop condition is hit"),
]

for num, step in steps:
    top = Inches(1.95) + (int(num)-1) * Inches(0.85)
    nb = s.shapes.add_shape(1, Inches(0.5), top, Inches(0.5), Inches(0.55))
    nb.fill.solid(); nb.fill.fore_color.rgb = ACCENT
    nb.line.fill.background()
    txbox(s, num, Inches(0.52), top, Inches(0.45), Inches(0.55),
          size=15, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    txbox(s, step, Inches(1.15), top + Inches(0.04), Inches(11.5), Inches(0.55),
          size=18, color=WHITE)

txbox(s, "The model owns the loop. You own the tools and the stop condition.",
      Inches(0.6), Inches(6.55), Inches(12), Inches(0.6),
      size=16, bold=True, color=ACCENT)


# ════════════════════════════════════════════════════════════════════
# SLIDE 13 — Core Principles
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "Three Principles for Reliable Agents",
      Inches(0.6), Inches(0.55), Inches(12), Inches(1.0),
      size=32, bold=True)
divider(s, Inches(1.65))

principles = [
    ("Simplicity",
     "Maintain Simple Design",
     "Add complexity only when measurable accuracy demands it. Most production systems don't need agents."),
    ("Transparency",
     "Show the Planning",
     "Surface the agent's reasoning and intermediate steps. You can't debug what you can't see."),
    ("Interface",
     "Engineer the ACI",
     "Agent-Computer Interface matters as much as UX. Tool docs, error messages, schemas — all part of the design."),
]

card_w = Inches(4.0)
for i, (label, title, desc) in enumerate(principles):
    left = Inches(0.4) + i * (card_w + Inches(0.16))
    card = s.shapes.add_shape(1, left, Inches(1.95), card_w, Inches(5.0))
    card.fill.solid(); card.fill.fore_color.rgb = TAG_BG
    card.line.color.rgb = ACCENT
    txbox(s, label.upper(), left + Inches(0.2), Inches(2.1),
          card_w - Inches(0.4), Inches(0.5), size=13, bold=True, color=ACCENT)
    txbox(s, title, left + Inches(0.2), Inches(2.6),
          card_w - Inches(0.4), Inches(0.9), size=20, bold=True, color=WHITE)
    txbox(s, desc, left + Inches(0.2), Inches(3.6),
          card_w - Inches(0.4), Inches(3.0), size=16, color=LIGHT_GREY)


# ════════════════════════════════════════════════════════════════════
# SLIDE 14 — Avoid / Instead
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "Common Agent Pitfalls",
      Inches(0.6), Inches(0.55), Inches(12), Inches(1.0),
      size=32, bold=True)
divider(s, Inches(1.65))

bad  = [
    "Reach for an agent on day one",
    "Hide the agent's reasoning from logs",
    "Hand the agent vague tool descriptions",
    "Skip the stop condition",
]
good = [
    "Try a single LLM call, then a workflow first",
    "Stream and log every step the agent takes",
    "Treat tool docs like API docs — be precise",
    "Cap iterations, time, or budget explicitly",
]

txbox(s, "AVOID", Inches(0.6), Inches(1.85), Inches(3), Inches(0.45),
      size=13, bold=True, color=RGBColor(0xFF, 0x55, 0x55))
txbox(s, "INSTEAD", Inches(7.3), Inches(1.85), Inches(3), Inches(0.45),
      size=13, bold=True, color=ACCENT)

for i, (b, g) in enumerate(zip(bad, good)):
    top = Inches(2.4) + i * Inches(1.05)
    bad_box = s.shapes.add_shape(1, Inches(0.5), top, Inches(6.3), Inches(0.85))
    bad_box.fill.solid(); bad_box.fill.fore_color.rgb = RGBColor(0x2A, 0x10, 0x10)
    bad_box.line.color.rgb = RGBColor(0xFF, 0x55, 0x55)
    txbox(s, f"X  {b}", Inches(0.7), top + Inches(0.18), Inches(6.0), Inches(0.6),
          size=16, color=RGBColor(0xFF, 0x88, 0x88))

    good_box = s.shapes.add_shape(1, Inches(7.0), top, Inches(5.8), Inches(0.85))
    good_box.fill.solid(); good_box.fill.fore_color.rgb = RGBColor(0x0D, 0x2A, 0x10)
    good_box.line.color.rgb = RGBColor(0x44, 0xBB, 0x44)
    txbox(s, f">  {g}", Inches(7.2), top + Inches(0.18), Inches(5.5), Inches(0.6),
          size=16, color=RGBColor(0x88, 0xDD, 0x88))


# ════════════════════════════════════════════════════════════════════
# SLIDE 15 — Closing
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "Agents Are a Tool.\nUse Them When the Task Demands It.",
      Inches(0.6), Inches(0.6), Inches(12), Inches(2.0),
      size=38, bold=True)
divider(s, Inches(2.8))

txbox(s, "Takeaways", Inches(0.6), Inches(3.0), Inches(6), Inches(0.7),
      size=24, bold=True, color=ACCENT)

pts = [
    "Start with the simplest thing that could work.",
    "Workflows for predictable tasks. Agents for open-ended ones.",
    "Pick a known pattern before inventing a new architecture.",
    "Make the agent's steps visible. Cap the loop. Invest in tool docs.",
    "Source: anthropic.com/engineering/building-effective-agents",
]
bullets(s, pts, Inches(0.6), Inches(3.7), Inches(12), Inches(3.5), size=19)


# ════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════
out = "/Users/amanparmar/Documents/AI-PM/Projects/next-leap-claude-code/02-presentation/Building-Effective-AI-Agents.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
