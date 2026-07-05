from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Colour palette
BG_DARK    = RGBColor(0x0D, 0x0D, 0x0D)
ACCENT     = RGBColor(0xD4, 0x7A, 0x21)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GREY   = RGBColor(0x88, 0x88, 0x88)
TAG_BG     = RGBColor(0x1E, 0x1E, 0x1E)
TAG_TEXT   = ACCENT
RED_BG     = RGBColor(0x2A, 0x10, 0x10)
RED_BORDER = RGBColor(0xFF, 0x55, 0x55)
RED_TEXT   = RGBColor(0xFF, 0x88, 0x88)
GREEN_BG   = RGBColor(0x0D, 0x2A, 0x10)
GREEN_BORDER = RGBColor(0x44, 0xBB, 0x44)
GREEN_TEXT = RGBColor(0x88, 0xDD, 0x88)
DARK_LINE  = RGBColor(0x33, 0x33, 0x33)
CODE_BG    = RGBColor(0x0A, 0x0A, 0x0A)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def txbox(slide, text, left, top, width, height,
          size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box

def accent_bar(slide, top=Inches(0.18), height=Inches(0.06)):
    bar = slide.shapes.add_shape(1, 0, top, SLIDE_W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

def bullets(slide, items, left, top, width, height,
            size=20, color=LIGHT_GREY, spacing_after=Pt(10)):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing_after
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color

def divider(slide, top, color=ACCENT):
    line = slide.shapes.add_shape(1, Inches(0.6), top, Inches(12.1), Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


# ============== SLIDE 1: Title ==============
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "CLAUDE CODE  ·  SKILLS EXPLAINER",
      Inches(0.6), Inches(0.7), Inches(12), Inches(0.5),
      size=13, color=MID_GREY)

txbox(s, "Claude Skills",
      Inches(0.6), Inches(1.4), Inches(11), Inches(1.1),
      size=54, bold=True, color=WHITE)

txbox(s, "Capabilities You Teach Once. Reuse Forever.",
      Inches(0.6), Inches(2.65), Inches(11), Inches(0.9),
      size=28, color=ACCENT)

divider(s, Inches(4.0))

txbox(s, "A skill is the difference between explaining a workflow 100 times\nand codifying it once. This deck covers what they are, how they work,\nand when to reach for one.",
      Inches(0.6), Inches(4.4), Inches(12), Inches(2.0),
      size=18, color=LIGHT_GREY)

txbox(s, "Sources: Anthropic Skills docs · Building Effective Agents · github.com/anthropics/skills",
      Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
      size=12, color=MID_GREY)


# ============== SLIDE 2: Agenda ==============
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "What This Deck Covers",
      Inches(0.6), Inches(0.55), Inches(11), Inches(0.9),
      size=34, bold=True)
divider(s, Inches(1.5))

rows = [
    ("01", "The problem skills solve"),
    ("02", "What a Claude Skill is"),
    ("03", "Anatomy of SKILL.md"),
    ("04", "The description field is the trigger"),
    ("05", "Progressive disclosure"),
    ("06", "Where skills live"),
    ("07", "Skills vs sub-agents vs CLAUDE.md"),
    ("08", "Three PM skill examples"),
    ("09", "How to build one"),
    ("10", "Anti-patterns + when to reach for a skill"),
]

row_h = Inches(0.49)
top_start = Inches(1.75)
for i, (num, title) in enumerate(rows):
    top = top_start + i * row_h
    if i % 2 == 0:
        shade = s.shapes.add_shape(1, Inches(0.5), top + Inches(0.03),
                                   Inches(12.3), row_h - Inches(0.04))
        shade.fill.solid()
        shade.fill.fore_color.rgb = TAG_BG
        shade.line.fill.background()
    txbox(s, num, Inches(0.65), top + Inches(0.08), Inches(0.8), row_h,
          size=14, color=ACCENT, bold=True)
    txbox(s, title, Inches(1.5), top + Inches(0.08), Inches(10.5), row_h,
          size=16, color=WHITE)


# ============== SLIDE 3: The problem ==============
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "Three Ways PMs Work With Claude Today",
      Inches(0.6), Inches(0.55), Inches(12), Inches(1.0),
      size=34, bold=True)
divider(s, Inches(1.6))

txbox(s, "All three break when the work gets repeatable.",
      Inches(0.6), Inches(1.75), Inches(11), Inches(0.5),
      size=16, color=MID_GREY)

problems = [
    ("01", "Copy-paste the same prompt every time",
     "You write the perfect prompt for feedback triage. Next week you write it again from memory. Slightly different. Slightly worse."),
    ("02", "Stuff everything into CLAUDE.md",
     "CLAUDE.md balloons past 500 lines. Every agent reads it on every call. Token cost goes up. Specific rules get lost in general context."),
    ("03", "Re-explain context every session",
     "You start each Claude session by pasting the same setup. Cohort context. Repo conventions. Your own preferences. Every. Single. Time."),
]

for i, (num, title, desc) in enumerate(problems):
    top = Inches(2.35) + i * Inches(1.55)
    bar = s.shapes.add_shape(1, Inches(0.5), top, Inches(12.3), Inches(1.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TAG_BG
    bar.line.color.rgb = ACCENT
    txbox(s, num, Inches(0.75), top + Inches(0.25), Inches(0.8), Inches(0.7),
          size=22, bold=True, color=ACCENT)
    txbox(s, title, Inches(1.7), top + Inches(0.2), Inches(10.5), Inches(0.5),
          size=18, bold=True, color=WHITE)
    txbox(s, desc, Inches(1.7), top + Inches(0.7), Inches(10.5), Inches(0.65),
          size=14, color=LIGHT_GREY)

txbox(s, "Skills are the layer that codifies repeatable work without bloating context.",
      Inches(0.6), Inches(7.05), Inches(12), Inches(0.4),
      size=15, bold=True, color=ACCENT)


# ============== SLIDE 4: What is a Claude Skill ==============
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "A Claude Skill Is a Reusable Capability With One Job",
      Inches(0.6), Inches(0.55), Inches(12), Inches(1.0),
      size=32, bold=True)
divider(s, Inches(1.65))

pts = [
    "Defined in a SKILL.md file with YAML frontmatter and a markdown body.",
    "Has a name (what it's called) and a description (when it should trigger).",
    "Any agent in the session can invoke it. No re-implementation per agent.",
    "Auto-loaded only when Claude decides the description matches the task.",
    "Lives in .claude/skills/ - shareable across projects, plugins, teammates.",
]
bullets(s, pts, Inches(0.6), Inches(1.95), Inches(7.6), Inches(4.5), size=19)

# Right panel - structure preview
panel = s.shapes.add_shape(1, Inches(8.5), Inches(1.85), Inches(4.3), Inches(4.4))
panel.fill.solid()
panel.fill.fore_color.rgb = TAG_BG
panel.line.color.rgb = ACCENT

txbox(s, "Skill folder structure",
      Inches(8.65), Inches(2.0), Inches(4.0), Inches(0.4),
      size=13, bold=True, color=ACCENT)

txbox(s, ".claude/skills/\n  feedback-triage/\n    SKILL.md\n    references/\n      rubric.md\n    scripts/\n      classify.py",
      Inches(8.65), Inches(2.55), Inches(4.0), Inches(3.5),
      size=14, color=ACCENT)

txbox(s, "A skill is one capability, in one folder, with one entry point.",
      Inches(0.6), Inches(6.6), Inches(12), Inches(0.5),
      size=15, bold=True, color=ACCENT)


# ============== SLIDE 5: Anatomy of SKILL.md ==============
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "Every SKILL.md = Frontmatter + Body",
      Inches(0.6), Inches(0.55), Inches(12), Inches(1.0),
      size=34, bold=True)
divider(s, Inches(1.65))

pts = [
    "Frontmatter (YAML, between --- markers):",
    "   name  - kebab-case identifier (required)",
    "   description - the TRIGGER (required, most important)",
    "   metadata - optional version, author, allowed-tools",
    "Body (markdown, below the frontmatter):",
    "   Instructions Claude follows when the skill fires.",
    "   Keep it under 500 lines. Move overflow to references/ and scripts/.",
]
bullets(s, pts, Inches(0.6), Inches(1.9), Inches(6.4), Inches(4.5), size=16, spacing_after=Pt(6))

# Right panel - actual SKILL.md example
panel = s.shapes.add_shape(1, Inches(7.1), Inches(1.85), Inches(5.8), Inches(4.4))
panel.fill.solid()
panel.fill.fore_color.rgb = CODE_BG
panel.line.color.rgb = ACCENT

txbox(s, "SKILL.md example",
      Inches(7.25), Inches(2.0), Inches(5.5), Inches(0.4),
      size=13, bold=True, color=ACCENT)

code = ("---\n"
        "name: feedback-triage\n"
        "description: Triage user feedback into\n"
        "  bug, feature, or churn signal. Use when\n"
        "  the user mentions feedback, tickets,\n"
        "  or customer complaints.\n"
        "---\n\n"
        "You classify incoming feedback into one\n"
        "of three buckets and return structured\n"
        "JSON for downstream agents.\n\n"
        "Steps:\n"
        "1. Read the feedback text.\n"
        "2. Apply the rubric (references/rubric.md).\n"
        "3. Return { category, confidence, evidence }.")
txbox(s, code, Inches(7.25), Inches(2.5), Inches(5.5), Inches(3.7),
      size=11, color=LIGHT_GREY)

txbox(s, "Frontmatter triggers the skill. Body tells Claude what to do.",
      Inches(0.6), Inches(6.6), Inches(12), Inches(0.5),
      size=15, bold=True, color=ACCENT)


# ============== SLIDE 6: Description is the trigger ==============
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "The Description Field Is the Trigger",
      Inches(0.6), Inches(0.55), Inches(12), Inches(1.0),
      size=34, bold=True)
divider(s, Inches(1.65))

txbox(s, "Claude doesn't load all skills. It picks them by reading descriptions.",
      Inches(0.6), Inches(1.85), Inches(12), Inches(0.5),
      size=16, color=MID_GREY)

# Avoid / Instead comparison
txbox(s, "AVOID", Inches(0.6), Inches(2.55), Inches(3), Inches(0.45),
      size=13, bold=True, color=RED_BORDER)
txbox(s, "INSTEAD", Inches(7.0), Inches(2.55), Inches(3), Inches(0.45),
      size=13, bold=True, color=ACCENT)

pairs = [
    ("Vague: \"Helps with PM stuff\"",
     "\"Triage user feedback into bug / feature / churn. Use when the user mentions feedback or tickets.\""),
    ("Abstract: \"Handles documents\"",
     "\"Review a PRD against rubric. Use when the user asks to review, score, or critique a spec.\""),
    ("Implementation-leaky: \"Calls the LLM with X\"",
     "\"Scan a competitor's positioning and threat angle. Use when the user mentions analyze, compare, or scan a competitor.\""),
]

for i, (bad, good) in enumerate(pairs):
    top = Inches(3.1) + i * Inches(1.15)
    bad_box = s.shapes.add_shape(1, Inches(0.5), top, Inches(6.2), Inches(0.95))
    bad_box.fill.solid()
    bad_box.fill.fore_color.rgb = RED_BG
    bad_box.line.color.rgb = RED_BORDER
    txbox(s, f"X  {bad}", Inches(0.65), top + Inches(0.15), Inches(5.9), Inches(0.7),
          size=14, color=RED_TEXT)

    good_box = s.shapes.add_shape(1, Inches(6.9), top, Inches(6.0), Inches(0.95))
    good_box.fill.solid()
    good_box.fill.fore_color.rgb = GREEN_BG
    good_box.line.color.rgb = GREEN_BORDER
    txbox(s, f"OK  {good}", Inches(7.05), top + Inches(0.12), Inches(5.7), Inches(0.75),
          size=12, color=GREEN_TEXT)

txbox(s, "Write the description like the search query a user would type to find this skill.",
      Inches(0.6), Inches(6.85), Inches(12), Inches(0.5),
      size=15, bold=True, color=ACCENT)


# ============== SLIDE 7: Progressive disclosure ==============
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "Progressive Disclosure: Load Less, Earlier",
      Inches(0.6), Inches(0.55), Inches(12), Inches(1.0),
      size=34, bold=True)
divider(s, Inches(1.65))

txbox(s, "Anthropic's pattern for keeping context lean as your skill library grows.",
      Inches(0.6), Inches(1.85), Inches(12), Inches(0.5),
      size=16, color=MID_GREY)

layers = [
    ("01", "Metadata", "name + description from every skill", "Always loaded into context, on every session. Costs almost nothing per skill.", Inches(2.65)),
    ("02", "Body", "SKILL.md instructions", "Loads ONLY when Claude triggers the skill based on the description match.", Inches(4.1)),
    ("03", "Resources", "references/ and scripts/", "Load on demand from inside the skill body. Use for rubrics, prompts, code that the skill calls.", Inches(5.55)),
]

for num, name, where, desc, top in layers:
    bar = s.shapes.add_shape(1, Inches(0.6), top, Inches(12.1), Inches(1.25))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TAG_BG
    bar.line.color.rgb = ACCENT
    txbox(s, num, Inches(0.8), top + Inches(0.2), Inches(0.7), Inches(0.9),
          size=24, bold=True, color=ACCENT)
    txbox(s, name, Inches(1.6), top + Inches(0.12), Inches(2.5), Inches(0.5),
          size=18, bold=True, color=WHITE)
    txbox(s, where, Inches(4.2), top + Inches(0.15), Inches(3.0), Inches(0.45),
          size=12, color=MID_GREY)
    txbox(s, desc, Inches(1.6), top + Inches(0.65), Inches(10.8), Inches(0.55),
          size=14, color=LIGHT_GREY)

txbox(s, "Same pattern lets you keep 50+ skills without bloating the context window.",
      Inches(0.6), Inches(7.05), Inches(12), Inches(0.4),
      size=15, bold=True, color=ACCENT)


# ============== SLIDE 8: Where skills live ==============
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "Three Locations - Same Skill File, Different Reach",
      Inches(0.6), Inches(0.55), Inches(12), Inches(1.0),
      size=32, bold=True)
divider(s, Inches(1.65))

cols = [
    ("GLOBAL", "~/.claude/skills/", "Every Claude Code session on your machine",
     "Your personal toolkit. Writing voice, daily routines, your-style triage."),
    ("PROJECT", "<repo>/.claude/skills/", "Only sessions inside this repo",
     "Repo-specific work. Conventions, codebase quirks, domain workflows."),
    ("PLUGIN", ".claude-plugin/skills/", "Anyone who installs your plugin",
     "Shippable to teammates. They install it; they get your capability."),
]

col_w = Inches(4.05)
for i, (label, path, scope, use) in enumerate(cols):
    left = Inches(0.4) + i * (col_w + Inches(0.16))
    card = s.shapes.add_shape(1, left, Inches(1.95), col_w, Inches(4.7))
    card.fill.solid()
    card.fill.fore_color.rgb = TAG_BG
    card.line.color.rgb = ACCENT

    txbox(s, label, left + Inches(0.18), Inches(2.1), col_w - Inches(0.3), Inches(0.5),
          size=14, bold=True, color=ACCENT)
    txbox(s, path, left + Inches(0.18), Inches(2.65), col_w - Inches(0.3), Inches(0.55),
          size=14, bold=True, color=WHITE)
    txbox(s, f"When: {scope}", left + Inches(0.18), Inches(3.3), col_w - Inches(0.3), Inches(0.9),
          size=12, color=MID_GREY)
    txbox(s, use, left + Inches(0.18), Inches(4.3), col_w - Inches(0.3), Inches(2.2),
          size=14, color=LIGHT_GREY)

txbox(s, "Move a skill from project to global when you want it everywhere. From global to plugin when you want to share it.",
      Inches(0.6), Inches(6.85), Inches(12), Inches(0.5),
      size=14, bold=True, color=ACCENT)


# ============== SLIDE 9: Skills vs sub-agents vs CLAUDE.md ==============
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "Skills vs Sub-agents vs CLAUDE.md",
      Inches(0.6), Inches(0.55), Inches(12), Inches(1.0),
      size=34, bold=True)
divider(s, Inches(1.65))

txbox(s, "The three are not interchangeable. Mixing them up causes most beginner pain.",
      Inches(0.6), Inches(1.85), Inches(12), Inches(0.5),
      size=15, color=MID_GREY)

cols = [
    ("SKILL", "Capability - HOW",
     ["What it does, not who does it",
      "Any agent can call it",
      "Loads only when triggered",
      "Example: feedback-triage skill"]),
    ("SUB-AGENT", "Persona - WHO",
     ["A named role with its own context",
      "Dispatched by an orchestrator",
      "Has its own tool budget",
      "Example: routing-agent that calls the feedback-triage skill"]),
    ("CLAUDE.md", "Context - WHAT IS THIS",
     ["Project background and conventions",
      "Always loaded for every agent",
      "One per project, one global",
      "Example: 'this is MeetFlow, a meeting AI'"]),
]

col_w = Inches(4.05)
for i, (label, subtitle, items) in enumerate(cols):
    left = Inches(0.4) + i * (col_w + Inches(0.16))
    card = s.shapes.add_shape(1, left, Inches(2.4), col_w, Inches(4.3))
    card.fill.solid()
    card.fill.fore_color.rgb = TAG_BG
    card.line.color.rgb = ACCENT

    txbox(s, label, left + Inches(0.18), Inches(2.55), col_w - Inches(0.3), Inches(0.45),
          size=14, bold=True, color=ACCENT)
    txbox(s, subtitle, left + Inches(0.18), Inches(3.05), col_w - Inches(0.3), Inches(0.55),
          size=17, bold=True, color=WHITE)
    bullets(s, items, left + Inches(0.1), Inches(3.75), col_w - Inches(0.2), Inches(2.7),
            size=13, color=LIGHT_GREY, spacing_after=Pt(6))

txbox(s, "Skill = how. Sub-agent = who. CLAUDE.md = what this project is.",
      Inches(0.6), Inches(6.95), Inches(12), Inches(0.4),
      size=15, bold=True, color=ACCENT)


# ============== SLIDE 10: PM Example 1 - feedback-triage ==============
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "PM Skill 1 - feedback-triage",
      Inches(0.6), Inches(0.55), Inches(11), Inches(1.0),
      size=34, bold=True)
divider(s, Inches(1.65))

pts = [
    "When it fires:",
    "   User mentions feedback, tickets, complaints, support requests",
    "What it does:",
    "   Reads the text, classifies as bug / feature / churn signal",
    "   Returns structured JSON for downstream agents",
    "Who calls it:",
    "   Monday routing agent, daily digest agent, weekly summary writer",
    "Why it's a skill, not a prompt:",
    "   Three different agents need the same classification logic",
]
bullets(s, pts, Inches(0.6), Inches(1.9), Inches(6.4), Inches(4.5), size=15, spacing_after=Pt(4))

# Right panel - SKILL.md preview
panel = s.shapes.add_shape(1, Inches(7.1), Inches(1.85), Inches(5.8), Inches(4.4))
panel.fill.solid()
panel.fill.fore_color.rgb = CODE_BG
panel.line.color.rgb = ACCENT

txbox(s, "SKILL.md",
      Inches(7.25), Inches(2.0), Inches(5.5), Inches(0.4),
      size=13, bold=True, color=ACCENT)

code = ("---\n"
        "name: feedback-triage\n"
        "description: Triage user feedback into bug,\n"
        "  feature, or churn. Use when the user\n"
        "  mentions feedback, tickets, or complaints.\n"
        "---\n\n"
        "Classify the incoming feedback. Return:\n"
        "  { category: bug | feature | churn,\n"
        "    confidence: 0-100,\n"
        "    evidence: \"quote from the text\" }\n\n"
        "If confidence < 70, return\n"
        "  { category: \"unclear\",\n"
        "    suggestion: \"escalate to human\" }")
txbox(s, code, Inches(7.25), Inches(2.5), Inches(5.5), Inches(3.7),
      size=11, color=LIGHT_GREY)

txbox(s, "One skill. Three agents. Consistent output. Update once.",
      Inches(0.6), Inches(6.55), Inches(12), Inches(0.5),
      size=15, bold=True, color=ACCENT)


# ============== SLIDE 11: PM Example 2 - prd-reviewer ==============
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "PM Skill 2 - prd-reviewer",
      Inches(0.6), Inches(0.55), Inches(11), Inches(1.0),
      size=34, bold=True)
divider(s, Inches(1.65))

pts = [
    "When it fires:",
    "   User asks to review, score, or critique a PRD or spec",
    "What it does:",
    "   Scores the doc 0-100 on clarity, acceptance criteria, edge cases",
    "   Returns issues + suggestions ranked by severity",
    "Who calls it:",
    "   Spec writer agent (Pattern 1 chain), eval-optimizer loop (Pattern 5)",
    "Why it's a skill:",
    "   The rubric is shared. Writer and reviewer both need it consistent",
]
bullets(s, pts, Inches(0.6), Inches(1.9), Inches(6.4), Inches(4.5), size=15, spacing_after=Pt(4))

# Right panel
panel = s.shapes.add_shape(1, Inches(7.1), Inches(1.85), Inches(5.8), Inches(4.4))
panel.fill.solid()
panel.fill.fore_color.rgb = CODE_BG
panel.line.color.rgb = ACCENT

txbox(s, "SKILL.md",
      Inches(7.25), Inches(2.0), Inches(5.5), Inches(0.4),
      size=13, bold=True, color=ACCENT)

code = ("---\n"
        "name: prd-reviewer\n"
        "description: Review a PRD against rubric.\n"
        "  Use when the user asks to review, score,\n"
        "  or critique a product spec or PRD.\n"
        "---\n\n"
        "Score 0-100 on:\n"
        "  - clarity of problem statement\n"
        "  - specificity of acceptance criteria\n"
        "  - edge case coverage\n"
        "  - dependency mention\n\n"
        "Return:\n"
        "  { score, issues: [...],\n"
        "    suggestions: [...],\n"
        "    verdict: pass | revise }")
txbox(s, code, Inches(7.25), Inches(2.5), Inches(5.5), Inches(3.7),
      size=11, color=LIGHT_GREY)

txbox(s, "Use the same rubric across spec generation, review, and audit.",
      Inches(0.6), Inches(6.55), Inches(12), Inches(0.5),
      size=15, bold=True, color=ACCENT)


# ============== SLIDE 12: PM Example 3 - competitive-scan ==============
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "PM Skill 3 - competitive-scan",
      Inches(0.6), Inches(0.55), Inches(11), Inches(1.0),
      size=34, bold=True)
divider(s, Inches(1.65))

pts = [
    "When it fires:",
    "   User asks to scan, analyze, or compare a competitor",
    "What it does:",
    "   Fetches positioning + pricing from public sources",
    "   Extracts threat angle and evidence URLs",
    "Who calls it:",
    "   Competitive research orchestrator, weekly intel digest, deal-prep agent",
    "Why it's a skill:",
    "   Same web-fetch + extraction pattern repeats across many agents",
]
bullets(s, pts, Inches(0.6), Inches(1.9), Inches(6.4), Inches(4.5), size=15, spacing_after=Pt(4))

# Right panel
panel = s.shapes.add_shape(1, Inches(7.1), Inches(1.85), Inches(5.8), Inches(4.4))
panel.fill.solid()
panel.fill.fore_color.rgb = CODE_BG
panel.line.color.rgb = ACCENT

txbox(s, "SKILL.md",
      Inches(7.25), Inches(2.0), Inches(5.5), Inches(0.4),
      size=13, bold=True, color=ACCENT)

code = ("---\n"
        "name: competitive-scan\n"
        "description: Scan a competitor's positioning\n"
        "  and threat angle from public sources.\n"
        "  Use when the user asks to analyze, compare,\n"
        "  or scan a competitor.\n"
        "---\n\n"
        "Read the competitor's homepage and pricing.\n"
        "Identify the threat angle (pricing, feature\n"
        "gap, GTM channel).\n\n"
        "Return:\n"
        "  { competitor, positioning,\n"
        "    threat_level, threat_summary,\n"
        "    evidence_urls: [...] }")
txbox(s, code, Inches(7.25), Inches(2.5), Inches(5.5), Inches(3.7),
      size=11, color=LIGHT_GREY)

txbox(s, "Composable: routing agents pass tickets, this skill does the work.",
      Inches(0.6), Inches(6.55), Inches(12), Inches(0.5),
      size=15, bold=True, color=ACCENT)


# ============== SLIDE 13: How to build a skill ==============
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "How to Build a Skill - Four Steps",
      Inches(0.6), Inches(0.55), Inches(12), Inches(1.0),
      size=34, bold=True)
divider(s, Inches(1.65))

steps = [
    ("01", "Scaffold with init_skill.py",
     "python3 scripts/init_skill.py <skill-name> --path .claude/skills/  - creates folder + SKILL.md with valid frontmatter"),
    ("02", "Write the description first",
     "It's the trigger. Write it like the search query a user would type. Include 'Use when ...' phrasing."),
    ("03", "Write the body next",
     "Instructions Claude follows. Keep under 500 lines. Move overflow into references/ and scripts/ subdirs."),
    ("04", "Test by triggering naturally",
     "Open Claude Code, paste a realistic request. If the skill doesn't fire, the description is the bug. Iterate."),
]

for i, (num, title, desc) in enumerate(steps):
    top = Inches(1.95) + i * Inches(1.18)
    bar = s.shapes.add_shape(1, Inches(0.5), top, Inches(12.3), Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TAG_BG
    bar.line.color.rgb = ACCENT
    txbox(s, num, Inches(0.75), top + Inches(0.18), Inches(0.7), Inches(0.7),
          size=22, bold=True, color=ACCENT)
    txbox(s, title, Inches(1.7), top + Inches(0.08), Inches(10.5), Inches(0.4),
          size=17, bold=True, color=WHITE)
    txbox(s, desc, Inches(1.7), top + Inches(0.55), Inches(10.5), Inches(0.5),
          size=13, color=LIGHT_GREY)

txbox(s, "The whole loop should take 15 minutes for your first real skill.",
      Inches(0.6), Inches(6.85), Inches(12), Inches(0.4),
      size=15, bold=True, color=ACCENT)


# ============== SLIDE 14: init_skill.py ==============
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "Use the Official Scaffolder, Not Blank Page",
      Inches(0.6), Inches(0.55), Inches(12), Inches(1.0),
      size=34, bold=True)
divider(s, Inches(1.65))

pts = [
    "Anthropic ships a skill-creator skill at github.com/anthropics/skills.",
    "Run init_skill.py to scaffold a valid SKILL.md with placeholders.",
    "Creates kebab-case folders, validates frontmatter, sets up subdirs.",
    "Saves you from common errors: missing fields, wrong filename, bad YAML.",
    "Also ships package_skill.py to validate and bundle skills for sharing.",
]
bullets(s, pts, Inches(0.6), Inches(1.9), Inches(6.5), Inches(4.0), size=16)

# Right panel - command
panel = s.shapes.add_shape(1, Inches(7.3), Inches(1.85), Inches(5.6), Inches(4.4))
panel.fill.solid()
panel.fill.fore_color.rgb = CODE_BG
panel.line.color.rgb = ACCENT

txbox(s, "Terminal",
      Inches(7.45), Inches(2.0), Inches(5.4), Inches(0.4),
      size=13, bold=True, color=ACCENT)

code = ("$ python3 scripts/init_skill.py \\\n"
        "    feedback-triage \\\n"
        "    --path .claude/skills/\n\n"
        "Created: .claude/skills/feedback-triage/\n"
        "  SKILL.md          (frontmatter ready)\n"
        "  references/       (empty, for overflow)\n"
        "  scripts/          (empty, for code)\n\n"
        "$ python3 scripts/package_skill.py \\\n"
        "    .claude/skills/feedback-triage\n"
        "Validated. Output: feedback-triage.zip")
txbox(s, code, Inches(7.45), Inches(2.5), Inches(5.4), Inches(3.7),
      size=12, color=LIGHT_GREY)

txbox(s, "Source: github.com/anthropics/skills/tree/main/skills/skill-creator",
      Inches(0.6), Inches(6.55), Inches(12), Inches(0.5),
      size=14, bold=True, color=ACCENT)


# ============== SLIDE 15: Anti-patterns ==============
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "What NOT to Do",
      Inches(0.6), Inches(0.55), Inches(11), Inches(0.9),
      size=34, bold=True)
divider(s, Inches(1.55))

txbox(s, "AVOID", Inches(0.6), Inches(1.65), Inches(3), Inches(0.45),
      size=13, bold=True, color=RED_BORDER)
txbox(s, "INSTEAD", Inches(7.0), Inches(1.65), Inches(3), Inches(0.45),
      size=13, bold=True, color=ACCENT)

pairs = [
    ("Vague description (\"PM helper\")",
     "Specific triggers (\"Use when user mentions PRD review\")"),
    ("SKILL.md body over 500 lines",
     "Split overflow into references/ and scripts/ subdirs"),
    ("Putting project context in a skill",
     "Project context goes in CLAUDE.md; skills are capabilities"),
    ("Building a skill for a single one-off prompt",
     "If one prompt works, just write the prompt - no skill needed"),
    ("Skill that mixes 4 different jobs",
     "One skill = one capability. Split into multiple skills."),
]

for i, (bad, good) in enumerate(pairs):
    top = Inches(2.2) + i * Inches(0.85)
    bad_box = s.shapes.add_shape(1, Inches(0.5), top, Inches(6.3), Inches(0.7))
    bad_box.fill.solid()
    bad_box.fill.fore_color.rgb = RED_BG
    bad_box.line.color.rgb = RED_BORDER
    txbox(s, f"X  {bad}", Inches(0.65), top + Inches(0.13), Inches(6.0), Inches(0.5),
          size=13, color=RED_TEXT)

    good_box = s.shapes.add_shape(1, Inches(7.0), top, Inches(5.8), Inches(0.7))
    good_box.fill.solid()
    good_box.fill.fore_color.rgb = GREEN_BG
    good_box.line.color.rgb = GREEN_BORDER
    txbox(s, f"OK  {good}", Inches(7.15), top + Inches(0.13), Inches(5.5), Inches(0.5),
          size=13, color=GREEN_TEXT)


# ============== SLIDE 16: Composition ==============
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "Skills Compose With the Rest of Claude Code",
      Inches(0.6), Inches(0.55), Inches(12), Inches(1.0),
      size=32, bold=True)
divider(s, Inches(1.65))

txbox(s, "A skill is one primitive. The full system has four. They snap together.",
      Inches(0.6), Inches(1.85), Inches(12), Inches(0.5),
      size=15, color=MID_GREY)

cards = [
    ("SKILLS",  "Capability",   "Called BY agents. Loaded when description matches."),
    ("AGENTS",  "Persona",      "Dispatch skills. Have their own scoped tools."),
    ("HOOKS",   "Deterministic glue", "Run at lifecycle events. Validate skill output."),
    ("MCP",     "External tools", "Provide tools that skills and agents can use."),
]

card_w = Inches(2.95)
card_top = Inches(2.55)
for i, (label, kind, desc) in enumerate(cards):
    left = Inches(0.4) + i * (card_w + Inches(0.18))
    card = s.shapes.add_shape(1, left, card_top, card_w, Inches(3.7))
    card.fill.solid()
    card.fill.fore_color.rgb = TAG_BG
    card.line.color.rgb = ACCENT
    txbox(s, label, left + Inches(0.18), card_top + Inches(0.2),
          card_w - Inches(0.3), Inches(0.55),
          size=16, bold=True, color=ACCENT)
    txbox(s, kind, left + Inches(0.18), card_top + Inches(0.8),
          card_w - Inches(0.3), Inches(0.55),
          size=14, bold=True, color=WHITE)
    txbox(s, desc, left + Inches(0.18), card_top + Inches(1.5),
          card_w - Inches(0.3), Inches(2.0),
          size=13, color=LIGHT_GREY)

txbox(s, "A plugin bundles all four into one installable unit your teammate can use.",
      Inches(0.6), Inches(6.55), Inches(12), Inches(0.5),
      size=15, bold=True, color=ACCENT)


# ============== SLIDE 17: When to reach for a skill ==============
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "When to Reach for a Skill",
      Inches(0.6), Inches(0.55), Inches(12), Inches(1.0),
      size=34, bold=True)
divider(s, Inches(1.65))

questions = [
    ("Can a single prompt do the job once?",
     "Then don't make a skill. Just write the prompt."),
    ("Do multiple agents or sessions need the same logic?",
     "Yes - make a skill. One source of truth."),
    ("Is the capability project-specific?",
     "Save it under <repo>/.claude/skills/."),
    ("Do you use it across many projects?",
     "Promote it to ~/.claude/skills/."),
    ("Do you want to share it with teammates?",
     "Bundle it into a plugin under .claude-plugin/skills/."),
]

for i, (q, a) in enumerate(questions):
    top = Inches(1.95) + i * Inches(0.92)
    bar = s.shapes.add_shape(1, Inches(0.5), top, Inches(12.3), Inches(0.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TAG_BG
    bar.line.color.rgb = ACCENT
    txbox(s, q, Inches(0.75), top + Inches(0.08), Inches(11.5), Inches(0.35),
          size=14, bold=True, color=WHITE)
    txbox(s, a, Inches(0.75), top + Inches(0.45), Inches(11.5), Inches(0.3),
          size=13, color=ACCENT)

txbox(s, "Start at the bottom rung. Promote upward only when reuse demands it.",
      Inches(0.6), Inches(6.75), Inches(12), Inches(0.5),
      size=15, bold=True, color=ACCENT)


# ============== SLIDE 18: Closing ==============
s = prs.slides.add_slide(BLANK)
bg(s); accent_bar(s)

txbox(s, "Skills Codify Repeatable Work\nWithout Bloating Context.",
      Inches(0.6), Inches(0.6), Inches(12), Inches(1.9),
      size=38, bold=True)
divider(s, Inches(2.6))

txbox(s, "Where to go next",
      Inches(0.6), Inches(2.8), Inches(6), Inches(0.6),
      size=22, bold=True, color=ACCENT)

links = [
    "Claude Code Skills docs - code.claude.com/docs/en/skills",
    "Anthropic skill-creator - github.com/anthropics/skills/tree/main/skills/skill-creator",
    "Building Effective Agents - anthropic.com/engineering/building-effective-agents",
    "Multi-Agent Research System - anthropic.com/engineering/multi-agent-research-system",
    "Sub-agents docs (skills + agents) - code.claude.com/docs/en/sub-agents",
]
bullets(s, links, Inches(0.6), Inches(3.5), Inches(12), Inches(2.5), size=16)

txbox(s, "Pick one workflow you do 3+ times a week. Make it a skill. That's it.",
      Inches(0.6), Inches(6.55), Inches(12), Inches(0.6),
      size=17, bold=True, color=ACCENT)


# SAVE
out = "/Users/amanparmar/Documents/AI-PM/Projects/next-leap-claude-code/outputs/Claude-Skills-Explainer.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
